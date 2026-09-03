"""Replace graphs in SIH2026-IDEA-Presentation-Submission.pptx and reports/START_TO_FINISH_PROJECT_REPORT.docx with updated CycML graphs without touching any user-edited text."""
import os
import shutil
import subprocess
import zipfile
from pathlib import Path
import pptx


def update_pptx():
    pptx_path = Path("SIH2026-IDEA-Presentation-Submission.pptx")
    if not pptx_path.exists():
        raise FileNotFoundError(f"File not found: {pptx_path}")

    prs = pptx.Presentation(pptx_path)
    print(f"Loaded presentation: {pptx_path} ({len(prs.slides)} slides)")

    # Mapping of slide index (0-based) to new figure path
    slide_graph_map = {
        2: "figures/slide3_technical_architecture_flowchart.png",  # Slide 3
        3: "figures/slide4_least_error_cyclones.png",               # Slide 4
        4: "figures/slide5_roadmap_infographic.png",                # Slide 5
        5: "figures/slide6_benchmark_comparison.png",               # Slide 6
    }

    for slide_idx, img_path in slide_graph_map.items():
        slide = prs.slides[slide_idx]
        img_bytes = Path(img_path).read_bytes()
        updated = False
        for shape in slide.shapes:
            # Find the main content graph (width > 6 inches or name Picture 17410, avoiding the top-right SIH logo)
            if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE and shape.left < 5000000:
                rId = shape._element.xpath('.//a:blip/@r:embed')[0]
                part = shape.part.related_part(rId)
                part._blob = img_bytes
                print(f"Slide {slide_idx+1}: Successfully replaced image part {rId} with {img_path} (new size: {len(img_bytes):,} bytes)")
                updated = True
                break
        if not updated:
            print(f"Warning: Main graph on Slide {slide_idx+1} not found.")

    prs.save(pptx_path)
    print(f"Successfully saved updated presentation to {pptx_path}")

    # Convert to PDF
    print("Converting updated presentation to PDF via LibreOffice...")
    cmd = ["libreoffice", "--headless", "--convert-to", "pdf", str(pptx_path), "--outdir", str(pptx_path.parent)]
    res = subprocess.run(cmd, capture_output=True, text=True)
    if res.returncode == 0:
        pdf_path = pptx_path.with_suffix(".pdf")
        print(f"Generated updated PDF: {pdf_path} ({pdf_path.stat().st_size / 1024 / 1024:.2f} MB)")
    else:
        print(f"LibreOffice PDF conversion error: {res.stderr}")


def update_docx():
    docx_path = Path("reports/START_TO_FINISH_PROJECT_REPORT.docx")
    if not docx_path.exists():
        raise FileNotFoundError(f"File not found: {docx_path}")

    # Map zip internal media files to new CycML figures
    media_map = {
        "word/media/image1.png": "figures/slide3_technical_architecture_flowchart.png",
        "word/media/image2.png": "figures/slide4_least_error_cyclones.png",
        "word/media/image3.png": "figures/slide6_benchmark_comparison.png",
        "word/media/image5.png": "figures/slide5_roadmap_infographic.png",
    }

    temp_docx = docx_path.with_suffix(".temp.docx")
    with zipfile.ZipFile(docx_path, "r") as zin:
        with zipfile.ZipFile(temp_docx, "w") as zout:
            for item in zin.infolist():
                if item.filename in media_map:
                    new_data = Path(media_map[item.filename]).read_bytes()
                    zout.writestr(item, new_data)
                    print(f"DOCX: Replaced {item.filename} with {media_map[item.filename]} ({len(new_data):,} bytes)")
                elif item.filename == "word/document.xml":
                    # Replace any remaining DeepCycloNet text with CycML
                    xml_content = zin.read(item.filename).decode("utf-8")
                    xml_content = xml_content.replace("DeepCycloNet", "CycML")
                    zout.writestr(item, xml_content.encode("utf-8"))
                    print("DOCX: Replaced any remaining DeepCycloNet text in document.xml with CycML")
                else:
                    zout.writestr(item, zin.read(item.filename))

    # Overwrite original
    shutil.move(str(temp_docx), str(docx_path))
    print(f"Successfully saved updated DOCX to {docx_path}")

    # Convert to PDF
    print("Converting updated DOCX to PDF via LibreOffice...")
    cmd = ["libreoffice", "--headless", "--convert-to", "pdf", str(docx_path), "--outdir", str(docx_path.parent)]
    res = subprocess.run(cmd, capture_output=True, text=True)
    if res.returncode == 0:
        pdf_path = docx_path.with_suffix(".pdf")
        print(f"Generated updated PDF: {pdf_path} ({pdf_path.stat().st_size / 1024 / 1024:.2f} MB)")
    else:
        print(f"LibreOffice PDF conversion error: {res.stderr}")


if __name__ == "__main__":
    print("=" * 70)
    print("UPDATING GRAPHS TO CycML")
    print("=" * 70)
    update_pptx()
    print("-" * 70)
    update_docx()
    print("=" * 70)
    print("DONE!")
