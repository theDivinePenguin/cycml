"""
Populates the official SIH 2026 PowerPoint Presentation (SIH2026-IDEA-Presentation-Format.pptx)
Strictly adheres to all SIH format rules, template styling, and required pointer headings.
Dual-column layout on Slide 2: Prediction graph on left, intuitive solution on right.
Cleans all stray bullet artifacts with explicit OpenXML buNone formatting.
"""

import shutil
from pathlib import Path
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml


def clear_bullet(p):
    """Disable automatic XML bullet so manual bullets render with 100% clean alignment."""
    pPr = p._p.get_or_add_pPr()
    for child in list(pPr):
        if any(tag in child.tag for tag in ['buFont', 'buChar', 'buAutoNum', 'buBlip']):
            pPr.remove(child)
    buNone = parse_xml('<a:buNone xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
    pPr.append(buNone)


def build_presentation():
    input_path = Path("SIH2026-IDEA-Presentation-Format.pptx")
    backup_path = Path("SIH2026-IDEA-Presentation-Format.backup.pptx")
    output_path = Path("SIH2026-IDEA-Presentation-Format.pptx")

    if not backup_path.exists():
        shutil.copyfile(input_path, backup_path)
        print(f"Created backup at {backup_path}")

    # Open fresh from backup to avoid compounding edits
    prs = pptx.Presentation(backup_path)

    FONT_FAMILY = "Arial"
    COLOR_PRIMARY = RGBColor(16, 54, 110)    # Navy Blue
    COLOR_BODY = RGBColor(33, 37, 41)        # Charcoal / Dark Grey

    # -------------------------------------------------------------
    # SLIDE 1: TITLE PAGE
    # -------------------------------------------------------------
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if shape.name == "TextBox 9" and shape.has_text_frame:
            tf = shape.text_frame
            tf.clear()
            
            lines = [
                ("Problem Statement ID – ", "26070", True),
                ("Problem Statement Title – ", "AI/ML-Driven Multi-Modal Spatio-Temporal Forecasting of Tropical Cyclone Intensity & Rapid Intensification", True),
                ("Theme – ", "Disaster Management / Space Technology & Smart Automation", True),
                ("PS Category – ", "Software", True),
                ("Team ID – ", "[Enter Registered Team ID on Portal]", False),
                ("Team Name – ", "[Enter Registered Team Name on Portal]", False),
            ]
            
            for idx, (label, val, highlight) in enumerate(lines):
                p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                clear_bullet(p)
                p.space_after = Pt(10)
                p.space_before = Pt(2)
                
                run1 = p.add_run()
                run1.text = label
                run1.font.name = FONT_FAMILY
                run1.font.size = Pt(15)
                run1.font.bold = True
                run1.font.color.rgb = COLOR_PRIMARY
                
                run2 = p.add_run()
                run2.text = val
                run2.font.name = FONT_FAMILY
                run2.font.size = Pt(15)
                run2.font.bold = highlight
                run2.font.color.rgb = COLOR_BODY

    # -------------------------------------------------------------
    # SLIDE 2: PROPOSED SOLUTION (Dual Column: Graph on Left, Solution on Right)
    # -------------------------------------------------------------
    slide2 = prs.slides[1]
    
    # 1. Update Title position so it does not touch the Team Oval
    for shape in slide2.shapes:
        if shape.name == "Title 1" and shape.has_text_frame:
            shape.left = Inches(1.85)
            shape.top = Inches(0.15)
            shape.width = Inches(8.8)
            shape.height = Inches(0.9)
            shape.text_frame.text = "DeepCycloNet: AI for Cyclone Tracking & Rapid Intensification"
            p_title = shape.text_frame.paragraphs[0]
            clear_bullet(p_title)
            p_title.font.name = FONT_FAMILY
            p_title.font.size = Pt(22)
            p_title.font.bold = True
            p_title.font.color.rgb = COLOR_PRIMARY

    # 2. Add Prediction Graph Image on the Left
    graph_img_path = "figures/lifecycle_lead_time_super_cyclone_phet.png"
    if Path(graph_img_path).exists():
        for s in list(slide2.shapes):
            if "Picture" in s.name and s.left < Inches(5):
                sp = s._element
                sp.getparent().remove(sp)

        img_left = Inches(0.6)
        img_top = Inches(1.5)
        img_width = Inches(5.3)
        img_height = Inches(3.3)
        slide2.shapes.add_picture(graph_img_path, img_left, img_top, width=img_width, height=img_height)

        # Add clean caption below graph
        caption_box = slide2.shapes.add_textbox(Inches(0.6), Inches(4.88), Inches(5.3), Inches(1.8))
        ctf = caption_box.text_frame
        ctf.word_wrap = True
        ctf.clear()
        
        cp1 = ctf.paragraphs[0]
        clear_bullet(cp1)
        cp1.text = "▲ Live Prototype Validation on Super Cyclone Phet (North Indian Ocean):"
        cp1.font.name = FONT_FAMILY
        cp1.font.size = Pt(10)
        cp1.font.bold = True
        cp1.font.color.rgb = COLOR_PRIMARY
        cp1.space_after = Pt(2)
        
        cp2 = ctf.add_paragraph()
        clear_bullet(cp2)
        cp2.text = "The red line shows actual observed storm wind speeds. The blue curve shows our AI's 24-hour advance forecast correctly anticipating the violent Rapid Intensification surge well before landfall."
        cp2.font.name = FONT_FAMILY
        cp2.font.size = Pt(9.5)
        cp2.font.color.rgb = COLOR_BODY

    # 3. Position Text Box on the Right with Plain-English, Non-Technical Content
    for shape in slide2.shapes:
        if shape.name == "TextBox 8" and shape.has_text_frame:
            shape.left = Inches(6.1)
            shape.top = Inches(1.4)
            shape.width = Inches(6.7)
            shape.height = Inches(5.2)

            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True

            # Header
            p_head = tf.paragraphs[0]
            clear_bullet(p_head)
            p_head.text = "❖ Proposed Solution (Describe your Idea/Solution/Prototype)"
            p_head.font.name = FONT_FAMILY
            p_head.font.size = Pt(14)
            p_head.font.bold = True
            p_head.font.color.rgb = COLOR_PRIMARY
            p_head.space_after = Pt(4)

            sections = [
                ("Detailed explanation of the proposed solution:", [
                    "An AI weather-intelligence system that watches a 12-hour video history of satellite cloud imagery (IR, Water Vapor & Visible) combined with real-time ocean temperature measurements to forecast tropical cyclone strength 24 hours ahead.",
                    "How it works: A Computer Vision network analyzes cloud rotation and eye formation, while an Ocean-Atmosphere model evaluates whether the sea below has enough heat energy to fuel sudden storm growth."
                ]),
                ("How it addresses the problem:", [
                    "Solves the Deadly Rapid Intensification Blindspot: Most cyclone casualties occur when weak storms unexpectedly explode into monster hurricanes right before landfall. Our AI detects this 24 hours early, giving disaster managers crucial lead time to evacuate coastlines.",
                    "Cuts Evacuation Costs & False Alarms: Distinguishes between storms that will intensify and those dying over cooler waters, avoiding unnecessary multi-crore evacuation expenses and public panic.",
                    "Reliable, Continuous Predictions: Anchors every prediction directly to the storm's current live wind speed, ensuring smooth forecasts without sudden unrealistic drop-offs."
                ]),
                ("Innovation and uniqueness of the solution:", [
                    "First Cloud + Ocean Dual AI: Combines satellite cloud visuals with ocean heat content maps, ensuring the AI checks both the storm's appearance and the oceanic fuel powering it.",
                    "Seamless 24/7 Day-and-Night Operation: Specially engineered to work continuously after sunset when visible optical satellite cameras go dark.",
                    "Proven on 1,265 Real Cyclones: Built and validated on 15+ years of real global storms with a working, interactive forecaster web dashboard."
                ])
            ]

            for heading, bullets in sections:
                p_ptr = tf.add_paragraph()
                clear_bullet(p_ptr)
                p_ptr.space_before = Pt(4)
                p_ptr.space_after = Pt(1)
                run_p = p_ptr.add_run()
                run_p.text = f"• {heading}"
                run_p.font.name = FONT_FAMILY
                run_p.font.size = Pt(11)
                run_p.font.bold = True
                run_p.font.color.rgb = COLOR_PRIMARY

                for b in bullets:
                    p_b = tf.add_paragraph()
                    clear_bullet(p_b)
                    p_b.space_before = Pt(1)
                    p_b.space_after = Pt(1)
                    run_b = p_b.add_run()
                    run_b.text = f"  - {b}"
                    run_b.font.name = FONT_FAMILY
                    run_b.font.size = Pt(9.5)
                    run_b.font.color.rgb = COLOR_BODY

    # Helper function for Slide 3, 4, 5, 6
    def format_text_slide(slide_idx, title_text, sections, font_size_b=10):
        s = prs.slides[slide_idx]
        for shape in s.shapes:
            if shape.name == "Title 1" and shape.has_text_frame:
                shape.left = Inches(1.85)
                shape.top = Inches(0.15)
                shape.width = Inches(8.8)
                shape.height = Inches(0.9)
                shape.text_frame.text = title_text
                p_t = shape.text_frame.paragraphs[0]
                clear_bullet(p_t)
                p_t.font.name = FONT_FAMILY
                p_t.font.size = Pt(22)
                p_t.font.bold = True
                p_t.font.color.rgb = COLOR_PRIMARY

            if shape.name == "TextBox 8" and shape.has_text_frame:
                shape.left = Inches(0.8)
                shape.top = Inches(1.35)
                shape.width = Inches(11.8)
                shape.height = Inches(5.3)

                tf = shape.text_frame
                tf.clear()
                tf.word_wrap = True

                for idx, (heading, bullets) in enumerate(sections):
                    p_ptr = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                    clear_bullet(p_ptr)
                    p_ptr.space_before = Pt(5)
                    p_ptr.space_after = Pt(2)
                    run_p = p_ptr.add_run()
                    run_p.text = f"• {heading}"
                    run_p.font.name = FONT_FAMILY
                    run_p.font.size = Pt(12)
                    run_p.font.bold = True
                    run_p.font.color.rgb = COLOR_PRIMARY

                    for b in bullets:
                        p_b = tf.add_paragraph()
                        clear_bullet(p_b)
                        p_b.space_before = Pt(1)
                        p_b.space_after = Pt(1)
                        run_b = p_b.add_run()
                        run_b.text = f"  - {b}"
                        run_b.font.name = FONT_FAMILY
                        run_b.font.size = Pt(font_size_b)
                        run_b.font.color.rgb = COLOR_BODY

    # -------------------------------------------------------------
    # SLIDE 3: TECHNICAL APPROACH
    # -------------------------------------------------------------
    s3_sections = [
        ("Technologies to be used (e.g. programming languages, frameworks, hardware):", [
            "Languages & Deep Learning: Python 3.11, PyTorch 2.4, TorchVision, CUDA 12, Mixed-Precision (FP16 GradScaler) for rapid high-throughput processing.",
            "Model Architecture: Frame-Shared ResNet-18 Feature Extractor, 2-Layer Temporal Transformer (8 attention heads, d_model=256), Gated Environmental Fusion MLP.",
            "Scientific Stack: NumPy, Pandas, HDF5 (h5py), Scikit-Learn, SciPy, Matplotlib, Seaborn, PyYAML.",
            "Deployment & Efficiency: Containerized with Docker / TorchScript for low-latency operational inference (<150 ms per storm cycle on GPU, <1.2s on standard CPU)."
        ]),
        ("Methodology and process for implementation (Flow Charts/Images/ working prototype):", [
            "1. Multi-Spectral Data Ingestion: Streams 5 consecutive geostationary satellite frames [t-12h, t-9h, t-6h, t-3h, t] covering Infrared (IR1 10.8µm), Water Vapor (WV 6.7µm), and Visible (VIS 0.65µm).",
            "2. Spatial Feature Extraction & Diurnal Gating: ResNet-18 extracts 512-d convective representations per frame; learned solar zenith gating dynamically handles daytime vs nighttime visible imagery.",
            "3. Temporal Sequence Modeling: Positional encoding + Temporal Transformer captures convective cloud rotation, eyewall contraction, and structural evolution over time.",
            "4. Environmental Fusion: Merges 6 physical ocean-atmosphere variables (SST, Ocean Heat Content, Shear, RH, Vmax, MSLP) into a unified multi-modal latent representation.",
            "5. Multi-Task Forecaster Output: Simultaneously outputs 24h RI Probability P(RI), 3-Class Trend (Weakening/Stable/Intensifying), and +6h/+12h/+24h continuous intensity guidance via an interactive web dashboard."
        ])
    ]
    format_text_slide(2, "TECHNICAL APPROACH", s3_sections, font_size_b=10)

    # -------------------------------------------------------------
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # -------------------------------------------------------------
    s4_sections = [
        ("Analysis of the feasibility of the idea:", [
            "Technical Feasibility: Fully functional, validated end-to-end prototype already developed and tested on 55,149 sequences across 1,265 real cyclones globally (North Indian Ocean, Pacific, Atlantic).",
            "Operational Latency: Forward inference takes <150 ms per sequence on modern GPU, easily fitting within the 30-minute operational dissemination window of IMD/JTWC.",
            "Data Availability: Operates on routine geostationary satellite feeds (INSAT-3D/3DR/3DS, GOES, Himawari) and GFS/SHIPS environmental analyses readily available in real time."
        ]),
        ("Potential challenges and risks:", [
            "Severe Class Imbalance: Rapid Intensification is rare (~6.8% prevalence in global cyclones), causing uncalibrated models to heavily under-predict RI.",
            "Missing Nighttime Visible Data: Visible sensors cannot capture cloud imagery at night, creating periodic 12-hour data gaps.",
            "Environmental Data Availability & Discrepancies: Global atmospheric analyses (GFS) are disseminated at 6h synoptic cycles, whereas satellite observations arrive every 3 hours."
        ]),
        ("Strategies for overcoming these challenges:", [
            "Cost-Sensitive Loss & Optimal Calibration: Weighted BCE loss (w_pos = 13.8) coupled with validation-optimized decision thresholds (tau = 0.141) achieves >63% RI recall and 0.868+ ROC-AUC.",
            "Learned Solar Zenith Gating: Dedicated gating layer smoothly downweights visible representations during nighttime passes, preventing sensor dropouts from degrading forecasts.",
            "Causal Zero-Lookahead Forward-Fill: Rigorous operational forward-fill from previous synoptic cycle with explicit audit logging (environment_age_hours), ensuring 100% causal mathematical safety."
        ])
    ]
    format_text_slide(3, "FEASIBILITY AND VIABILITY", s4_sections, font_size_b=9.5)

    # -------------------------------------------------------------
    # SLIDE 5: IMPACT AND BENEFITS
    # -------------------------------------------------------------
    s5_sections = [
        ("Potential impact on the target audience:", [
            "Direct Operational Beneficiaries: India Meteorological Department (IMD), National Disaster Management Authority (NDMA), State Disaster Management Authorities (SDMAs - Odisha, Gujarat, AP, WB, Tamil Nadu), Indian Coast Guard, and NDRF.",
            "24-Hour Actionable Evacuation Window: Provides coastal district administrators a life-saving 24-hour lead time to execute targeted evacuations, recall deep-sea fishermen, and reinforce embankments before extreme winds make roads impassable.",
            "Objective Probabilistic Decision Support: Replaces qualitative, subjective Dvorak human estimates with objective, calibrated, and reproducible AI confidence scores."
        ]),
        ("Benefits of the solution (social, economic, environmental, etc.):", [
            "Social (Preserving Human Lives): Directly mitigates mass casualties caused by unpredicted explosive offshore intensification (e.g., VSCS Nargis 2008: 138,000+ deaths; Cyclone Ockhi 2017: 350+ fishermen lost).",
            "Economic (Optimized Resource Allocation): Mass coastal evacuations cost state governments ₹50–100+ Crores per event. Accurate RI probability and trend classification prevents costly false alarm evacuations while ensuring full mobilization when danger is severe.",
            "Infrastructure & Maritime Protection: Gives port authorities (Paradip, Kandla, Vizag), naval assets, and offshore installations (ONGC platforms) vital lead time to secure heavy cranes, vessels, and fuel terminals against destructive storm surges."
        ])
    ]
    format_text_slide(4, "IMPACT AND BENEFITS", s5_sections, font_size_b=10)

    # -------------------------------------------------------------
    # SLIDE 6: RESEARCH AND REFERENCES
    # -------------------------------------------------------------
    s6_sections = [
        ("Details / Links of the reference and research work:", [
            "[1] Chen, B. F., Chen, B., Lin, K. T., & Elsberry, R. L. (2020). 'TCIR: A Tropical Cyclone Information Dataset for Deep Learning Applications.' IEEE Transactions on Geoscience and Remote Sensing (TGRS), 58(4), 2838–2847. DOI: 10.1109/TGRS.2019.2942474.",
            "[2] DeMaria, M., Mainelli, M., Shay, L. K., Knaff, J. A., & Kaplan, J. (2005). 'Further Improvements to the Statistical Hurricane Intensity Prediction Scheme (SHIPS).' Weather and Forecasting, 20(4), 530–544.",
            "[3] Kaplan, J., & DeMaria, M. (2003). 'Large-Scale Characteristics of Rapidly Intensifying Tropical Cyclones in the North Atlantic Basin.' Weather and Forecasting, 18(6), 1093–1108.",
            "[4] Vaswani, A., Shazeer, N., Parmar, N., Uszkoreit, J., Jones, L., Gomez, A. N., Kaiser, Ł., & Polosukhin, I. (2017). 'Attention Is All You Need.' Advances in Neural Information Processing Systems (NeurIPS 2017), 30, 5998–6008.",
            "[5] He, K., Zhang, X., Ren, S., & Sun, J. (2016). 'Deep Residual Learning for Image Recognition.' Proceedings of the IEEE Conference on Computer Vision and Pattern Recognition (CVPR 2016), 770–778.",
            "[6] CIRA / NOAA RAMMB Tropical Cyclone Diagnostics & SHIPS Database: https://rammb-data.cira.colostate.edu/ships/",
            "[7] India Meteorological Department (IMD) Cyclone Best Track Archives & INSAT Imagery: https://rsmcnewdelhi.imd.gov.in/",
            "[8] Project Prototype Repository & Real-Time Forecaster Demo: https://github.com/theDivinePenguin/cycml"
        ])
    ]
    format_text_slide(5, "RESEARCH  AND REFERENCES", s6_sections, font_size_b=9.5)

    # Save filled presentation
    prs.save(output_path)
    print(f"Successfully populated and saved: {output_path}")

    named_output = Path("SIH2026-DeepCycloNet-Submission.pptx")
    prs.save(named_output)
    print(f"Also saved copy at: {named_output}")


if __name__ == "__main__":
    build_presentation()
